import os
import sys
from pathlib import Path

# Add parent directory to path to allow imports when run directly
sys.path.insert(0, str(Path(__file__).parent.parent))

from mcp.server.fastmcp import FastMCP
from mcp.types import TextContent
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
import asyncio
from server_mcp.tools_arithmetic import (
    number_list_to_sum,
    add_numbers,
    calculate_difference,
    number_list_to_product,
    calculate_division,
    strings_to_chars_to_int as local_strings_to_chars_to_int,
    int_list_to_exponential_values,
    fibonacci_numbers as local_fibonacci_numbers,
    calculate_factorial,
    calculate_permutation,
    calculate_combination,
    calculate_salary_for_id,
    calculate_salary_for_name,
    calculate_percentage,
    # New arithmetic functions
    calculate_absolute_value,
    calculate_modulo,
    calculate_floor_division,
    calculate_ceiling,
    calculate_floor,
    calculate_round,
    calculate_gcd,
    calculate_lcm,
    is_prime,
    find_prime_factors,
    calculate_average,
    find_max,
    find_min,
    calculate_square,
    calculate_square_root,
    calculate_cube,
    calculate_cube_root,
    convert_to_fraction,
    calculate_reciprocal,
)

# Import logical reasoning tools
from server_mcp.tools_logical import (
    evaluate_logical_and,
    evaluate_logical_or,
    evaluate_logical_not,
    evaluate_implication,
    evaluate_biconditional,
    evaluate_xor,
    solve_syllogism,
    count_true_values,
    majority_vote,
    evaluate_complex_expression,
)

# Import algebra tools
from server_mcp.tools_algebra import (
    solve_linear_equation,
    solve_quadratic_equation,
    evaluate_polynomial,
    solve_system_2x2,
    calculate_power,
    calculate_nth_root,
    expand_binomial,
    calculate_arithmetic_sequence_sum,
    calculate_geometric_sequence_sum,
    find_arithmetic_sequence_term,
    find_geometric_sequence_term,
    simplify_ratio,
    parse_linear_equation,
    parse_quadratic_equation,
)

# Import geometry tools
from server_mcp.tools_geometry import (
    calculate_circle_area,
    calculate_circle_circumference,
    calculate_rectangle_area,
    calculate_rectangle_perimeter,
    calculate_triangle_area,
    calculate_triangle_area_heron,
    calculate_sphere_volume,
    calculate_sphere_surface_area,
    calculate_cylinder_volume,
    calculate_cylinder_surface_area,
    calculate_cone_volume,
    calculate_cube_volume,
    calculate_cube_surface_area,
    calculate_rectangular_prism_volume,
    calculate_distance_2d,
    calculate_distance_3d,
    calculate_pythagorean_theorem,
    calculate_pythagorean_leg,
    calculate_chord_length,
    calculate_trapezoid_area,
    calculate_parallelogram_area,
)

# Import statistics tools
from server_mcp.tools_statistics import (
    calculate_mean,
    calculate_median,
    calculate_mode,
    calculate_range,
    calculate_variance,
    calculate_standard_deviation,
    calculate_percentile,
    calculate_quartiles,
    calculate_interquartile_range,
    calculate_z_score,
    calculate_correlation_coefficient,
    calculate_linear_regression,
    calculate_factorial_stat,
    calculate_combinations_stat,
    calculate_probability_union,
    calculate_probability_complement,
)
from server_mcp.models import (
    DrawRectangleInput,
    AddTextInput,
    SendGmailInput,
    NumberListInput, NumberListOutput,
    TwoNumberInput, TwoNumberOutput,
    SingleNumberInput, SingleNumberOutput,
    SingleIntInput, TwoIntInput,
    RoundInput, PrimeCheckOutput, PrimeFactorsOutput,
    FractionOutput, FractionInput,
    PercentageInput, PercentageOutput,
    StringToCharsInput, StringToCharsOutput,
    ExponentialInput, ExponentialOutput,
    FibonacciInput, FibonacciOutput,
    FactorialInput, FactorialOutput,
    PermutationInput, PermutationOutput,
    CombinationInput, CombinationOutput,
    EmployeeIdInput, EmployeeNameInput, SalaryOutput,
    FallbackInput, FallbackOutput,
    # Logical models
    BooleanListInput, BooleanOutput,
    TwoBooleanInput, SingleBooleanInput,
    LogicalExpressionInput, IntOutput,
    # Algebra models
    LinearEquationInput, LinearEquationOutput,
    QuadraticEquationInput, QuadraticEquationOutput,
    PolynomialInput, System2x2Input, System2x2Output,
    PowerInput, RootInput, BinomialExpansionInput,
    FloatListOutput, ArithmeticSequenceInput, GeometricSequenceInput,
    RatioInput, RatioOutput,
    # Geometry models
    RadiusInput, RectangleInput,
    TriangleBaseHeightInput, TriangleThreeSidesInput,
    CylinderInput, CubeInput, RectangularPrismInput,
    Point2DInput, Point3DInput, PythagoreanInput, PythagoreanLegInput, ChordInput, TrapezoidInput,
    # Statistics models
    StatNumberListInput, StatFloatListOutput,
    VarianceInput, PercentileInput, QuartilesOutput,
    ZScoreInput, CorrelationInput, RegressionOutput,
    ProbabilityUnionInput, ProbabilityInput,
)
import logging
from datetime import datetime
import traceback
import smtplib
from email.mime.text import MIMEText
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
log_dir = "logs"
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, f"mcp_server_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log")

# Define constant for PowerPoint filename
PPTX_FILENAME = 'presentation.pptx'

logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Instantiate an MCP server client
mcp = FastMCP("Calculator")

# DEFINE TOOLS
@mcp.tool()
async def close_powerpoint() -> dict:
    """Close PowerPoint"""
    try:
        logger.info("Calling close_powerpoint()")
        # Suppress all output from taskkill
        try:
            # Use subprocess with output suppressed
            proc = await asyncio.create_subprocess_shell(
                'taskkill /F /IM POWERPNT.EXE',
                stdout=asyncio.subprocess.DEVNULL,
                stderr=asyncio.subprocess.DEVNULL
            )
            await proc.communicate()
        except Exception as e:
            logger.warning(f"taskkill failed or PowerPoint not running: {e}")
        await asyncio.sleep(2)
        logger.info("PowerPoint closed successfully")
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint closed successfully"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in close_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error closing PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def open_powerpoint() -> dict:
    """Open a new PowerPoint presentation"""
    try:
        logger.info("Calling open_powerpoint()")
        await close_powerpoint()
        await asyncio.sleep(3)
        prs = Presentation()
        prs.slide_layouts[0]
        filename = PPTX_FILENAME
        prs.save(filename)
        await asyncio.sleep(5)
        # Open the presentation, suppressing output
        try:
            # On Windows, os.startfile does not print to stdout, but just in case:
            # Use subprocess with output suppressed for other OSes if needed
            os.startfile(filename)
        except Exception as e:
            logger.warning(f"os.startfile failed: {e}")
        await asyncio.sleep(10)
        logger.info("PowerPoint opened successfully with a new presentation")
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint opened successfully with a new presentation"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in open_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def draw_rectangle(input: DrawRectangleInput) -> dict:
    """Draw a rectangle in the first slide of PowerPoint"""
    try:
        logger.info(f"Drawing rectangle with validated parameters: ({input.x1},{input.y1}) to ({input.x2},{input.y2})")
        
        # Wait before modifying the presentation
        await asyncio.sleep(2)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        await asyncio.sleep(2)
        prs = Presentation(PPTX_FILENAME)
        slide = prs.slides[0]
        
        # Store existing text boxes
        slide = prs.slides[0]
        
        # Store existing text boxes
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text_boxes.append((text, left, top, width, height))
        
        # Clear existing shapes except text boxes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                sp = shape._element
                sp.getparent().remove(sp)
        
        # Convert coordinates to inches
        left = Inches(input.x1)
        top = Inches(input.y1)
        width = Inches(input.x2 - input.x1)
        height = Inches(input.y2 - input.y1)
        
        logger.debug(f"Rectangle dimensions - left={left}, top={top}, width={width}, height={height}")
        
        # Add rectangle
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        
        # Make the rectangle more visible
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
        shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        prs.save(PPTX_FILENAME)
        await asyncio.sleep(2)
        
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(5)
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(5)
        
        logger.info(f"Rectangle drawn successfully from ({input.x1},{input.y1}) to ({input.x2},{input.y2})")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Rectangle drawn successfully from ({input.x1},{input.y1}) to ({input.x2},{input.y2})"
                )
            ]
        }
            
    except Exception as e:
        error_msg = f"Error in draw_rectangle: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return {"content": [TextContent(type="text", text=error_msg)]}

@mcp.tool()
async def add_text_in_powerpoint(input: AddTextInput) -> dict:
    """Add text to the first slide of PowerPoint"""
    try:
        logger.info(f"Received text to add: {input.text}")
        logger.debug(f"Text length: {len(input.text)}")
        logger.debug(f"Text contains newlines: {'\\n' in input.text}")
        
        # Wait before adding text
        await asyncio.sleep(5)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        prs = Presentation(PPTX_FILENAME)
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        # Match the rectangle position from draw_rectangle
        left = Inches(2.2)  # Slightly more than rectangle left for margin
        top = Inches(2.5)   # Centered vertically in rectangle
        width = Inches(4.6) # Slightly less than rectangle width for margin
        height = Inches(2)  # Enough height for text
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True  # Enable word wrap
        text_frame.vertical_anchor = 1  # Middle vertical alignment
        
        # Split text into lines
        lines = input.text.split('\n')
        logger.debug(f"Number of lines: {len(lines)}")
        logger.debug(f"Lines to add: {lines}")
        
        # Add each line as a separate paragraph
        for i, line in enumerate(lines):
            if line.strip():  # Only add non-empty lines
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.alignment = 1  # Center align the text
                
                # Format the text
                run = p.runs[0]
                if "Final Result:" in line:
                    run.font.size = Pt(32)  # Header size
                    run.font.bold = True
                else:
                    run.font.size = Pt(28)  # Value size
                    run.font.bold = True
                
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        prs.save(PPTX_FILENAME)
        await asyncio.sleep(5)
        
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(10)
        # Reopen PowerPoint
        os.startfile(PPTX_FILENAME)
        await asyncio.sleep(10)
        
        logger.info(f"Text added successfully: {input.text}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Text added successfully: {input.text}"
                )
            ]
        }
    except Exception as e:
        logger.error(f"Error in add_text_in_powerpoint: {str(e)}")
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def send_gmail(input: SendGmailInput) -> dict:
    """Send an email with the specified content via Gmail"""
    try:
        logger.info(f"Calling send_gmail(content: {input.content[:50]}...)")
        
        # Retrieve Gmail credentials and recipient from .env
        gmail_address = os.getenv("GMAIL_ADDRESS")
        gmail_app_password = os.getenv("GMAIL_APP_PASSWORD")
        recipient_email = os.getenv("RECIPIENT_EMAIL")
        
        if not all([gmail_address, gmail_app_password, recipient_email]):
            error_msg = "Missing GMAIL_ADDRESS, GMAIL_APP_PASSWORD, or RECIPIENT_EMAIL in .env file"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        # Validate Gmail address and recipient email format
        if not (gmail_address.endswith('@gmail.com') and '@' in recipient_email):
            error_msg = f"Invalid email format: GMAIL_ADDRESS={gmail_address}, RECIPIENT_EMAIL={recipient_email}"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        # Create the email message
        msg = MIMEText(input.content)
        msg['Subject'] = 'Math Agent Result'
        msg['From'] = gmail_address
        msg['To'] = recipient_email
        
        # Connect to Gmail's SMTP server
        try:
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
                logger.debug("Connecting to Gmail SMTP server (smtp.gmail.com:465)")
                server.login(gmail_address, gmail_app_password)
                logger.debug(f"Logged in as {gmail_address}")
                server.sendmail(gmail_address, recipient_email, msg.as_string())
                logger.info(f"Email sent successfully to {recipient_email}")
        except smtplib.SMTPAuthenticationError as e:
            error_msg = f"SMTP Authentication failed: {str(e)}. Ensure GMAIL_APP_PASSWORD is correct and 2-Step Verification is enabled."
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        except Exception as e:
            error_msg = f"Failed to send email: {str(e)}"
            logger.error(error_msg)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=error_msg
                    )
                ]
            }
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Email sent successfully to {recipient_email}"
                )
            ]
        }
    except Exception as e:
        error_msg = f"Error in send_gmail: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return {
            "content": [
                TextContent(
                    type="text",
                    text=error_msg
                )
            ]
        }

# TOOLS WRAPPING FUNCTIONS FROM tools.py (with Pydantic validation)
@mcp.tool()
def t_number_list_to_sum(input: NumberListInput) -> NumberListOutput:
    """Sum numbers in a list"""
    logger.info(f"Calling t_number_list_to_sum with {len(input.numbers)} numbers")
    result = number_list_to_sum(input.numbers)
    return NumberListOutput(result=result)

@mcp.tool()
def t_add(input: TwoNumberInput) -> TwoNumberOutput:
    """Add two numbers"""
    logger.info(f"Calling t_add({input.a} + {input.b})")
    result = add_numbers(input.a, input.b)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_calculate_difference(input: TwoNumberInput) -> TwoNumberOutput:
    """Difference between two numbers"""
    logger.info(f"Calling t_calculate_difference({input.a}, {input.b})")
    result = calculate_difference(input.a, input.b)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_number_list_to_product(input: NumberListInput) -> NumberListOutput:
    """Product of numbers in a list"""
    logger.info(f"Calling t_number_list_to_product with {len(input.numbers)} numbers")
    result = number_list_to_product(input.numbers)
    return NumberListOutput(result=result)

@mcp.tool()
def t_calculate_division(input: TwoNumberInput) -> TwoNumberOutput:
    """Division of two numbers"""
    logger.info(f"Calling t_calculate_division({input.a}, {input.b})")
    result = calculate_division(input.a, input.b)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_strings_to_chars_to_int(input: StringToCharsInput) -> StringToCharsOutput:
    """ASCII values of characters"""
    logger.info(f"Calling t_strings_to_chars_to_int('{input.text}')")
    ascii_values = local_strings_to_chars_to_int(input.text)
    return StringToCharsOutput(ascii_values=ascii_values)

@mcp.tool()
def t_int_list_to_exponential_values(input: ExponentialInput) -> ExponentialOutput:
    """Exponential of list elements"""
    logger.info(f"Calling t_int_list_to_exponential_values with {len(input.numbers)} numbers")
    values = int_list_to_exponential_values(input.numbers)
    return ExponentialOutput(values=values)

@mcp.tool()
def t_fibonacci_numbers(input: FibonacciInput) -> FibonacciOutput:
    """First n Fibonacci numbers"""
    logger.info(f"Calling t_fibonacci_numbers(n={input.n})")
    sequence = local_fibonacci_numbers(input.n)
    return FibonacciOutput(sequence=sequence)

@mcp.tool()
def t_calculate_factorial(input: FactorialInput) -> FactorialOutput:
    """List of factorials up to n-1"""
    logger.info(f"Calling t_calculate_factorial(n={input.n})")
    factorials = calculate_factorial(input.n)
    return FactorialOutput(factorials=factorials)

@mcp.tool()
def t_calculate_permutation(input: PermutationInput) -> PermutationOutput:
    """Permutation nPr"""
    logger.info(f"Calling t_calculate_permutation(n={input.n}, r={input.r})")
    result = calculate_permutation(input.n, input.r)
    return PermutationOutput(result=result)

@mcp.tool()
def t_calculate_combination(input: CombinationInput) -> CombinationOutput:
    """Combination nCr"""
    logger.info(f"Calling t_calculate_combination(n={input.n}, r={input.r})")
    result = calculate_combination(input.n, input.r)
    return CombinationOutput(result=result)

@mcp.tool()
def t_calculate_salary_for_id(input: EmployeeIdInput) -> SalaryOutput:
    """Salary by employee id"""
    logger.info(f"Calling t_calculate_salary_for_id(emp_id={input.emp_id})")
    salary = calculate_salary_for_id(input.emp_id)
    return SalaryOutput(salary=salary, found=salary is not None)

@mcp.tool()
def t_calculate_salary_for_name(input: EmployeeNameInput) -> SalaryOutput:
    """Salary by employee name"""
    logger.info(f"Calling t_calculate_salary_for_name(emp_name='{input.emp_name}')")
    salary = calculate_salary_for_name(input.emp_name)
    return SalaryOutput(salary=salary, found=salary is not None)

@mcp.tool()
def t_calculate_percentage(input: PercentageInput) -> PercentageOutput:
    """Calculate percentage of a number"""
    logger.info(f"Calling t_calculate_percentage(percent={input.percent}, number={input.number})")
    result = calculate_percentage(input.percent, input.number)
    return PercentageOutput(result=result)

@mcp.tool()
def t_absolute_value(input: SingleNumberInput) -> SingleNumberOutput:
    """Calculate absolute value of a number"""
    logger.info(f"Calling t_absolute_value({input.value})")
    result = calculate_absolute_value(input.value)
    return SingleNumberOutput(result=result)

@mcp.tool()
def t_modulo(input: TwoNumberInput) -> TwoNumberOutput:
    """Calculate modulo (remainder)"""
    logger.info(f"Calling t_modulo({input.a}, {input.b})")
    result = calculate_modulo(input.a, input.b)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_floor_division(input: TwoNumberInput) -> TwoNumberOutput:
    """Calculate floor division (integer division)"""
    logger.info(f"Calling t_floor_division({input.a}, {input.b})")
    result = calculate_floor_division(input.a, input.b)
    return TwoNumberOutput(result=float(result))

@mcp.tool()
def t_ceiling(input: SingleNumberInput) -> SingleNumberOutput:
    """Round number up to nearest integer"""
    logger.info(f"Calling t_ceiling({input.value})")
    result = calculate_ceiling(input.value)
    return SingleNumberOutput(result=float(result))

@mcp.tool()
def t_floor(input: SingleNumberInput) -> SingleNumberOutput:
    """Round number down to nearest integer"""
    logger.info(f"Calling t_floor({input.value})")
    result = calculate_floor(input.value)
    return SingleNumberOutput(result=float(result))

@mcp.tool()
def t_round(input: RoundInput) -> SingleNumberOutput:
    """Round number to specified decimal places"""
    logger.info(f"Calling t_round({input.number}, {input.decimals} decimals)")
    result = calculate_round(input.number, input.decimals)
    return SingleNumberOutput(result=result)

@mcp.tool()
def t_gcd(input: TwoIntInput) -> IntOutput:
    """Calculate Greatest Common Divisor (GCD)"""
    logger.info(f"Calling t_gcd({input.a}, {input.b})")
    result = calculate_gcd(input.a, input.b)
    return IntOutput(result=result)

@mcp.tool()
def t_lcm(input: TwoIntInput) -> IntOutput:
    """Calculate Least Common Multiple (LCM)"""
    logger.info(f"Calling t_lcm({input.a}, {input.b})")
    result = calculate_lcm(input.a, input.b)
    return IntOutput(result=result)

@mcp.tool()
def t_is_prime(input: SingleIntInput) -> PrimeCheckOutput:
    """Check if a number is prime"""
    logger.info(f"Calling t_is_prime({input.n})")
    result = is_prime(input.n)
    return PrimeCheckOutput(result=result, number=input.n)

@mcp.tool()
def t_prime_factors(input: SingleIntInput) -> PrimeFactorsOutput:
    """Find all prime factors of a number"""
    logger.info(f"Calling t_prime_factors({input.n})")
    factors = find_prime_factors(input.n)
    return PrimeFactorsOutput(factors=factors)

@mcp.tool()
def t_average(input: NumberListInput) -> TwoNumberOutput:
    """Calculate average of numbers"""
    logger.info(f"Calling t_average with {len(input.numbers)} numbers")
    result = calculate_average(input.numbers)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_max(input: NumberListInput) -> TwoNumberOutput:
    """Find maximum value in list"""
    logger.info(f"Calling t_max with {len(input.numbers)} numbers")
    result = find_max(input.numbers)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_min(input: NumberListInput) -> TwoNumberOutput:
    """Find minimum value in list"""
    logger.info(f"Calling t_min with {len(input.numbers)} numbers")
    result = find_min(input.numbers)
    return TwoNumberOutput(result=result)

@mcp.tool()
def t_square(input: SingleNumberInput) -> SingleNumberOutput:
    """Calculate square of a number"""
    logger.info(f"Calling t_square({input.value})")
    result = calculate_square(input.value)
    return SingleNumberOutput(result=result)

@mcp.tool()
def t_square_root(input: SingleNumberInput) -> SingleNumberOutput:
    """Calculate square root of a number"""
    logger.info(f"Calling t_square_root({input.value})")
    result = calculate_square_root(input.value)
    return SingleNumberOutput(result=result)

@mcp.tool()
def t_cube(input: SingleNumberInput) -> SingleNumberOutput:
    """Calculate cube of a number"""
    logger.info(f"Calling t_cube({input.value})")
    result = calculate_cube(input.value)
    return SingleNumberOutput(result=result)

@mcp.tool()
def t_cube_root(input: SingleNumberInput) -> SingleNumberOutput:
    """Calculate cube root of a number"""
    logger.info(f"Calling t_cube_root({input.value})")
    result = calculate_cube_root(input.value)
    return SingleNumberOutput(result=result)

@mcp.tool()
def t_to_fraction(input: FractionInput) -> FractionOutput:
    """Convert decimal to fraction"""
    logger.info(f"Calling t_to_fraction({input.decimal})")
    numerator, denominator = convert_to_fraction(input.decimal, input.max_denominator)
    return FractionOutput(numerator=numerator, denominator=denominator)

@mcp.tool()
def t_reciprocal(input: SingleNumberInput) -> SingleNumberOutput:
    """Calculate reciprocal (1/x) of a number"""
    logger.info(f"Calling t_reciprocal({input.value})")
    result = calculate_reciprocal(input.value)
    return SingleNumberOutput(result=result)

@mcp.tool()
def fallback_reasoning(input: FallbackInput) -> FallbackOutput:
    """Fallback reasoning step when the agent is uncertain or a tool fails"""
    logger.info(f"Calling fallback_reasoning: {input.description}")
    message = f"Fallback invoked: {input.description}"
    return FallbackOutput(message=message)


# ============================================
# LOGICAL REASONING TOOLS
# ============================================

@mcp.tool()
def t_logical_and(input: BooleanListInput) -> BooleanOutput:
    """Evaluate logical AND of boolean values"""
    logger.info(f"Calling t_logical_and with {len(input.values)} values")
    result = evaluate_logical_and(input.values)
    return BooleanOutput(result=result)


@mcp.tool()
def t_logical_or(input: BooleanListInput) -> BooleanOutput:
    """Evaluate logical OR of boolean values"""
    logger.info(f"Calling t_logical_or with {len(input.values)} values")
    result = evaluate_logical_or(input.values)
    return BooleanOutput(result=result)


@mcp.tool()
def t_logical_not(input: SingleBooleanInput) -> BooleanOutput:
    """Evaluate logical NOT"""
    logger.info(f"Calling t_logical_not({input.value})")
    result = evaluate_logical_not(input.value)
    return BooleanOutput(result=result)


@mcp.tool()
def t_implication(input: TwoBooleanInput) -> BooleanOutput:
    """Evaluate logical implication (premise → conclusion)"""
    logger.info(f"Calling t_implication({input.a} → {input.b})")
    result = evaluate_implication(input.a, input.b)
    return BooleanOutput(result=result)


@mcp.tool()
def t_biconditional(input: TwoBooleanInput) -> BooleanOutput:
    """Evaluate biconditional (a ↔ b)"""
    logger.info(f"Calling t_biconditional({input.a} ↔ {input.b})")
    result = evaluate_biconditional(input.a, input.b)
    return BooleanOutput(result=result)


@mcp.tool()
def t_xor(input: TwoBooleanInput) -> BooleanOutput:
    """Evaluate exclusive OR (XOR)"""
    logger.info(f"Calling t_xor({input.a} ⊕ {input.b})")
    result = evaluate_xor(input.a, input.b)
    return BooleanOutput(result=result)


@mcp.tool()
def t_syllogism(input: TwoBooleanInput) -> BooleanOutput:
    """Solve syllogism using modus ponens"""
    logger.info(f"Calling t_syllogism({input.a}, {input.b})")
    result = solve_syllogism(input.a, input.b)
    return BooleanOutput(result=result)


@mcp.tool()
def t_count_true(input: BooleanListInput) -> IntOutput:
    """Count number of true values"""
    logger.info(f"Calling t_count_true with {len(input.values)} values")
    result = count_true_values(input.values)
    return IntOutput(result=result)


@mcp.tool()
def t_majority_vote(input: BooleanListInput) -> BooleanOutput:
    """Determine majority vote from boolean values"""
    logger.info(f"Calling t_majority_vote with {len(input.values)} values")
    result = majority_vote(input.values)
    return BooleanOutput(result=result)


@mcp.tool()
def t_complex_expression(input: LogicalExpressionInput) -> BooleanOutput:
    """Evaluate complex logical expression"""
    logger.info(f"Calling t_complex_expression: {input.expression}")
    result = evaluate_complex_expression(input.expression, input.variables)
    return BooleanOutput(result=result)


# ============================================
# ALGEBRA TOOLS
# ============================================

@mcp.tool()
def t_solve_linear(input: LinearEquationInput) -> LinearEquationOutput:
    """Solve linear equation ax + b = 0 or from equation string like 'x + 4 = 5'"""
    # Check if equation_string is provided
    if input.equation_string:
        logger.info(f"Calling t_solve_linear with equation string: {input.equation_string}")
        try:
            a, b = parse_linear_equation(input.equation_string)
            logger.info(f"Parsed to: {a}x + {b} = 0")
        except Exception as e:
            logger.error(f"Failed to parse equation string: {e}")
            return LinearEquationOutput(solution=None)
    else:
        a = input.a
        b = input.b
        logger.info(f"Calling t_solve_linear({a}x + {b} = 0)")
    
    solution = solve_linear_equation(a, b)
    return LinearEquationOutput(solution=solution)


@mcp.tool()
def t_solve_quadratic(input: QuadraticEquationInput) -> QuadraticEquationOutput:
    """Solve quadratic equation ax² + bx + c = 0 or from equation string like 'x^2 - 5x + 6 = 0'"""
    # Check if equation_string is provided
    if input.equation_string:
        logger.info(f"Calling t_solve_quadratic with equation string: {input.equation_string}")
        try:
            a, b, c = parse_quadratic_equation(input.equation_string)
            logger.info(f"Parsed to: {a}x² + {b}x + {c} = 0")
        except Exception as e:
            logger.error(f"Failed to parse equation string: {e}")
            return QuadraticEquationOutput(solutions=[])
    else:
        a = input.a
        b = input.b
        c = input.c
        logger.info(f"Calling t_solve_quadratic({a}x² + {b}x + {c} = 0)")
    
    solutions = solve_quadratic_equation(a, b, c)
    return QuadraticEquationOutput(solutions=solutions)


@mcp.tool()
def t_evaluate_polynomial(input: PolynomialInput) -> TwoNumberOutput:
    """Evaluate polynomial at given x value"""
    logger.info(f"Calling t_evaluate_polynomial at x={input.x}")
    result = evaluate_polynomial(input.coefficients, input.x)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_solve_system_2x2(input: System2x2Input) -> System2x2Output:
    """Solve 2x2 system of linear equations"""
    logger.info("Calling t_solve_system_2x2")
    solution = solve_system_2x2(input.a1, input.b1, input.c1, input.a2, input.b2, input.c2)
    if solution:
        return System2x2Output(x=solution[0], y=solution[1], has_solution=True)
    return System2x2Output(x=None, y=None, has_solution=False)


@mcp.tool()
def t_power(input: PowerInput) -> TwoNumberOutput:
    """Calculate base raised to exponent"""
    logger.info(f"Calling t_power({input.base}^{input.exponent})")
    result = calculate_power(input.base, input.exponent)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_nth_root(input: RootInput) -> TwoNumberOutput:
    """Calculate nth root of a number"""
    logger.info(f"Calling t_nth_root({input.n}th root of {input.number})")
    result = calculate_nth_root(input.number, input.n)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_binomial_expansion(input: BinomialExpansionInput) -> FloatListOutput:
    """Expand binomial (a + b)^n"""
    logger.info(f"Calling t_binomial_expansion({input.a} + {input.b})^{input.n}")
    values = expand_binomial(input.a, input.b, input.n)
    return FloatListOutput(values=values)


@mcp.tool()
def t_arithmetic_sum(input: ArithmeticSequenceInput) -> TwoNumberOutput:
    """Calculate sum of arithmetic sequence"""
    logger.info("Calling t_arithmetic_sum")
    result = calculate_arithmetic_sequence_sum(input.first, input.common_diff, input.n)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_geometric_sum(input: GeometricSequenceInput) -> TwoNumberOutput:
    """Calculate sum of geometric sequence"""
    logger.info("Calling t_geometric_sum")
    result = calculate_geometric_sequence_sum(input.first, input.ratio, input.n)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_arithmetic_term(input: ArithmeticSequenceInput) -> TwoNumberOutput:
    """Find nth term of arithmetic sequence"""
    logger.info("Calling t_arithmetic_term")
    result = find_arithmetic_sequence_term(input.first, input.common_diff, input.n)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_geometric_term(input: GeometricSequenceInput) -> TwoNumberOutput:
    """Find nth term of geometric sequence"""
    logger.info("Calling t_geometric_term")
    result = find_geometric_sequence_term(input.first, input.ratio, input.n)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_simplify_ratio(input: RatioInput) -> RatioOutput:
    """Simplify ratio to lowest terms"""
    logger.info(f"Calling t_simplify_ratio({input.a}:{input.b})")
    a, b = simplify_ratio(input.a, input.b)
    return RatioOutput(simplified_a=a, simplified_b=b)


# ============================================
# GEOMETRY TOOLS
# ============================================

@mcp.tool()
def t_circle_area(input: RadiusInput) -> TwoNumberOutput:
    """Calculate area of a circle"""
    logger.info(f"Calling t_circle_area(r={input.radius})")
    result = calculate_circle_area(input.radius)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_circle_circumference(input: RadiusInput) -> TwoNumberOutput:
    """Calculate circumference of a circle"""
    logger.info(f"Calling t_circle_circumference(r={input.radius})")
    result = calculate_circle_circumference(input.radius)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_rectangle_area(input: RectangleInput) -> TwoNumberOutput:
    """Calculate area of a rectangle"""
    logger.info(f"Calling t_rectangle_area({input.length}x{input.width})")
    result = calculate_rectangle_area(input.length, input.width)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_rectangle_perimeter(input: RectangleInput) -> TwoNumberOutput:
    """Calculate perimeter of a rectangle"""
    logger.info(f"Calling t_rectangle_perimeter({input.length}x{input.width})")
    result = calculate_rectangle_perimeter(input.length, input.width)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_triangle_area(input: TriangleBaseHeightInput) -> TwoNumberOutput:
    """Calculate area of a triangle using base and height"""
    logger.info(f"Calling t_triangle_area(base={input.base}, height={input.height})")
    result = calculate_triangle_area(input.base, input.height)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_triangle_area_heron(input: TriangleThreeSidesInput) -> TwoNumberOutput:
    """Calculate area of a triangle using Heron's formula"""
    logger.info(f"Calling t_triangle_area_heron({input.a}, {input.b}, {input.c})")
    result = calculate_triangle_area_heron(input.a, input.b, input.c)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_sphere_volume(input: RadiusInput) -> TwoNumberOutput:
    """Calculate volume of a sphere"""
    logger.info(f"Calling t_sphere_volume(r={input.radius})")
    result = calculate_sphere_volume(input.radius)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_sphere_surface_area(input: RadiusInput) -> TwoNumberOutput:
    """Calculate surface area of a sphere"""
    logger.info(f"Calling t_sphere_surface_area(r={input.radius})")
    result = calculate_sphere_surface_area(input.radius)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_cylinder_volume(input: CylinderInput) -> TwoNumberOutput:
    """Calculate volume of a cylinder"""
    logger.info(f"Calling t_cylinder_volume(r={input.radius}, h={input.height})")
    result = calculate_cylinder_volume(input.radius, input.height)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_cylinder_surface_area(input: CylinderInput) -> TwoNumberOutput:
    """Calculate surface area of a cylinder"""
    logger.info(f"Calling t_cylinder_surface_area(r={input.radius}, h={input.height})")
    result = calculate_cylinder_surface_area(input.radius, input.height)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_cone_volume(input: CylinderInput) -> TwoNumberOutput:
    """Calculate volume of a cone"""
    logger.info(f"Calling t_cone_volume(r={input.radius}, h={input.height})")
    result = calculate_cone_volume(input.radius, input.height)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_cube_volume(input: CubeInput) -> TwoNumberOutput:
    """Calculate volume of a cube"""
    logger.info(f"Calling t_cube_volume(side={input.side})")
    result = calculate_cube_volume(input.side)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_cube_surface_area(input: CubeInput) -> TwoNumberOutput:
    """Calculate surface area of a cube"""
    logger.info(f"Calling t_cube_surface_area(side={input.side})")
    result = calculate_cube_surface_area(input.side)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_prism_volume(input: RectangularPrismInput) -> TwoNumberOutput:
    """Calculate volume of a rectangular prism"""
    logger.info(f"Calling t_prism_volume({input.length}x{input.width}x{input.height})")
    result = calculate_rectangular_prism_volume(input.length, input.width, input.height)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_distance_2d(input: Point2DInput) -> TwoNumberOutput:
    """Calculate distance between two 2D points"""
    logger.info("Calling t_distance_2d")
    result = calculate_distance_2d(input.x1, input.y1, input.x2, input.y2)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_distance_3d(input: Point3DInput) -> TwoNumberOutput:
    """Calculate distance between two 3D points"""
    logger.info("Calling t_distance_3d")
    result = calculate_distance_3d(input.x1, input.y1, input.z1, input.x2, input.y2, input.z2)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_pythagorean(input: PythagoreanInput) -> TwoNumberOutput:
    """Calculate hypotenuse using Pythagorean theorem given two legs"""
    logger.info(f"Calling t_pythagorean({input.a}, {input.b})")
    result = calculate_pythagorean_theorem(input.a, input.b)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_pythagorean_leg(input: PythagoreanLegInput) -> TwoNumberOutput:
    """Calculate unknown leg using Pythagorean theorem given one leg and hypotenuse"""
    logger.info(f"Calling t_pythagorean_leg(known_leg={input.known_leg}, hypotenuse={input.hypotenuse})")
    result = calculate_pythagorean_leg(input.known_leg, input.hypotenuse)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_chord_length(input: ChordInput) -> TwoNumberOutput:
    """Calculate length of a chord in a circle given radius and distance from center"""
    logger.info(f"Calling t_chord_length(radius={input.radius}, distance={input.distance_from_center})")
    result = calculate_chord_length(input.radius, input.distance_from_center)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_trapezoid_area(input: TrapezoidInput) -> TwoNumberOutput:
    """Calculate area of a trapezoid"""
    logger.info("Calling t_trapezoid_area")
    result = calculate_trapezoid_area(input.base1, input.base2, input.height)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_parallelogram_area(input: TriangleBaseHeightInput) -> TwoNumberOutput:
    """Calculate area of a parallelogram"""
    logger.info(f"Calling t_parallelogram_area(base={input.base}, height={input.height})")
    result = calculate_parallelogram_area(input.base, input.height)
    return TwoNumberOutput(result=result)


# ============================================
# STATISTICS TOOLS
# ============================================

@mcp.tool()
def t_mean(input: StatNumberListInput) -> TwoNumberOutput:
    """Calculate arithmetic mean (average)"""
    logger.info(f"Calling t_mean with {len(input.numbers)} numbers")
    result = calculate_mean(input.numbers)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_median(input: StatNumberListInput) -> TwoNumberOutput:
    """Calculate median"""
    logger.info(f"Calling t_median with {len(input.numbers)} numbers")
    result = calculate_median(input.numbers)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_mode(input: StatNumberListInput) -> StatFloatListOutput:
    """Calculate mode(s)"""
    logger.info(f"Calling t_mode with {len(input.numbers)} numbers")
    values = calculate_mode(input.numbers)
    return StatFloatListOutput(values=values)


@mcp.tool()
def t_range(input: StatNumberListInput) -> TwoNumberOutput:
    """Calculate range (max - min)"""
    logger.info(f"Calling t_range with {len(input.numbers)} numbers")
    result = calculate_range(input.numbers)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_variance(input: VarianceInput) -> TwoNumberOutput:
    """Calculate variance"""
    logger.info(f"Calling t_variance (sample={input.sample})")
    result = calculate_variance(input.numbers, input.sample)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_std_deviation(input: VarianceInput) -> TwoNumberOutput:
    """Calculate standard deviation"""
    logger.info(f"Calling t_std_deviation (sample={input.sample})")
    result = calculate_standard_deviation(input.numbers, input.sample)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_percentile(input: PercentileInput) -> TwoNumberOutput:
    """Calculate percentile"""
    logger.info(f"Calling t_percentile({input.percentile}th percentile)")
    result = calculate_percentile(input.numbers, input.percentile)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_quartiles(input: StatNumberListInput) -> QuartilesOutput:
    """Calculate quartiles (Q1, Q2, Q3)"""
    logger.info("Calling t_quartiles")
    q1, q2, q3 = calculate_quartiles(input.numbers)
    return QuartilesOutput(q1=q1, q2=q2, q3=q3)


@mcp.tool()
def t_iqr(input: StatNumberListInput) -> TwoNumberOutput:
    """Calculate interquartile range (IQR)"""
    logger.info("Calling t_iqr")
    result = calculate_interquartile_range(input.numbers)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_z_score(input: ZScoreInput) -> TwoNumberOutput:
    """Calculate z-score"""
    logger.info(f"Calling t_z_score({input.value})")
    result = calculate_z_score(input.value, input.mean, input.std_dev)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_correlation(input: CorrelationInput) -> TwoNumberOutput:
    """Calculate Pearson correlation coefficient"""
    logger.info("Calling t_correlation")
    result = calculate_correlation_coefficient(input.x_values, input.y_values)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_linear_regression(input: CorrelationInput) -> RegressionOutput:
    """Calculate linear regression (y = mx + b)"""
    logger.info("Calling t_linear_regression")
    slope, intercept = calculate_linear_regression(input.x_values, input.y_values)
    return RegressionOutput(slope=slope, intercept=intercept)


@mcp.tool()
def t_factorial_stat(input: FactorialInput) -> IntOutput:
    """Calculate factorial for statistics"""
    logger.info(f"Calling t_factorial_stat({input.n}!)")
    result = calculate_factorial_stat(input.n)
    return IntOutput(result=result)


@mcp.tool()
def t_combinations_stat(input: CombinationInput) -> IntOutput:
    """Calculate combinations C(n,r) for statistics"""
    logger.info(f"Calling t_combinations_stat(C({input.n},{input.r}))")
    result = calculate_combinations_stat(input.n, input.r)
    return IntOutput(result=result)


@mcp.tool()
def t_probability_union(input: ProbabilityUnionInput) -> TwoNumberOutput:
    """Calculate probability union P(A ∪ B)"""
    logger.info("Calling t_probability_union")
    result = calculate_probability_union(input.p_a, input.p_b, input.p_both)
    return TwoNumberOutput(result=result)


@mcp.tool()
def t_probability_complement(input: ProbabilityInput) -> TwoNumberOutput:
    """Calculate probability complement P(A')"""
    logger.info("Calling t_probability_complement")
    result = calculate_probability_complement(input.p)
    return TwoNumberOutput(result=result)


if __name__ == "__main__":
    logger.info("Starting the MCP server")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        logger.info("Running MCP server in dev mode without transport")
        mcp.run()  # Run without transport for dev server
    else:
        logger.info("Running MCP server with stdio transport")
        mcp.run(transport="stdio")  # Run with stdio for direct execution